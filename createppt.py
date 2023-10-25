from pptx import Presentation

# Create a PowerPoint presentation object
prs = Presentation()

# Add a title slide
title_slide_layout = prs.slide_layouts[0]  # 0 is the layout for the title slide
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
title.text = "My Presentation Title"

# Add a content slide
bullet_slide_layout = prs.slide_layouts[1]  # 1 is the layout for a content slide
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = "Slide 1 Title"
tf = body_shape.text_frame
tf.text = "This is a sample bullet point slide."

# Add some bullet points
p = tf.add_paragraph()
p.text = "Bullet Point 1"
p.level = 0  # Level 0 is the top level (main bullet point)
p = tf.add_paragraph()
p.text = "Bullet Point 2"
p.level = 1  # Level 1 is a sub-bullet

# Save the presentation to a file
prs.save("sample.pptx")
